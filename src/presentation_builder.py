import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from models import SlideContent, VisualType

class PresentationBuilder:
    def __init__(self , template_path: str = None):
        
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
            
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)
        self._apply_theme()
        
    def _apply_theme(self):
        # Define a modern color palette and background for all slides
        for slide in self.prs.slides:
            self._set_slide_background(slide)
        
    def _set_slide_background(self, slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(245, 245, 245)  # White background
        
    # def _set_slide_background_color(self, slide=None):
    #     if slide is None:
    #         for slide in self.prs.slides: 
    #             fill = slide.background.fill
    #             fill.solid()
    #             fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
    #     else: 
    #         fill = slide.background.fill
    #         fill.solid()
    #         fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
        
    def add_title_slide(self, company_name: str, subtitle: str):
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide)
        
        title = slide.shapes.title
        title.text = company_name
        title.text_frame.paragraphs[0].font.size = Pt(60) 
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 70, 122)  # Blue
        title.text_frame.alignment = PP_ALIGN.CENTER
        
        subtitle_placeholder = slide.placeholders[1]
        subtitle_placeholder.text = subtitle
        subtitle_placeholder.text_frame.paragraphs[0].font.size = Pt(32) 
        subtitle_placeholder.text_frame.paragraphs[0].font.italic = True
        subtitle_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)  # Gray
        subtitle_placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def add_content_slide(self, content: SlideContent):
        blank_layout = self.prs.slide_layouts[6] if len(self.prs.slide_layouts) > 6 else self.prs.slide_layouts[-1]
        slide = self.prs.slides.add_slide(blank_layout)
        self._set_slide_background(slide)
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(1))
        tf_title = title_box.text_frame
        tf_title.text = content.title.upper()
        p_title = tf_title.paragraphs[0]
        p_title.font.bold = True
        p_title.font.size = Pt(44)
        p_title.font.color.rgb = RGBColor(0, 70, 122)  
        p_title.alignment = PP_ALIGN.CENTER
            
        # Add text content
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(5))
        tf_content = content_box.text_frame
        tf_content.word_wrap = True
        for line in content.content.split('\n'):
            p = tf_content.add_paragraph()
            p.text = line
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(50, 50, 50)
            p.space_after = Pt(10)
        
        # Add visual if present
        if content.visual_type and content.visual_data:
            visual_left = Inches(9)
            visual_top = Inches(1.5)
            visual_width = Inches(6)
            if content.visual_type == VisualType.CHART:
                self._add_chart(slide, content.visual_data, visual_left, visual_top, visual_width)
            elif content.visual_type == VisualType.TIMELINE:
                self._add_timeline(slide, content.visual_data, visual_left, visual_top, visual_width)
            elif content.visual_type == VisualType.IMAGE:
                self._add_image(slide, content.visual_data, visual_left, visual_top, visual_width)
                
    def _add_chart(self, slide, visual_data, visual_left, visual_top, visual_width):
        if 'chart' in visual_data and visual_data['chart'] is not None: 
            chart_filepath = "temp_chart.png" 
            chart = visual_data['chart']
            chart.update_layout(
                width=800,
                height=600,
                margin=dict(l=20, r=20, t=40, b=20),
                title_font_size=24,
                plot_bgcolor='rgba(0,0,0,0)',
                paper_bgcolor='rgba(0,0,0,0)'
            )
            
            chart.write_image(chart_filepath)
            try:
                slide.shapes.add_picture(chart_filepath, visual_left, visual_top, width=visual_width)
            except Exception as e:
                print(f"Error adding chart to slide: {e}") 
            finally:
                if os.path.exists(chart_filepath):
                    os.remove(chart_filepath)  

    def _add_timeline(self, slide, visual_data, visual_left, visual_top, visual_width):
        if 'chart' in visual_data and visual_data['chart'] is not None: 
            timeline_filepath = "temp_timeline.png"
            timeline_chart = visual_data['chart']
            timeline_chart.write_image(timeline_filepath)
            try:
                slide.shapes.add_picture(timeline_filepath, visual_left, visual_top, width=visual_width)
            except Exception as e:
                print(f"Error adding timeline to slide: {e}") 
            finally:
                if os.path.exists(timeline_filepath):
                    os.remove(timeline_filepath)
      
    def _add_image(self, slide, visual_data, left, top, width):
        # shape = slide.shapes.add_shape(
        #     MSO_SHAPE.RECTANGLE,
        #     left=left, top=top, width=width, height=height
        # )
        # shape.fill.background() 
        # shape.line.color.rgb = RGBColor(84, 110, 122)  # Gray
        # shape.line.width = Pt(1)
        # text_frame = shape.text_frame
        # text_frame.text = "Image Placeholder"
        # text_frame.paragraphs[0].font.size = Pt(16)
        # text_frame.paragraphs[0].font.color.rgb = RGBColor(84, 110, 122)  # Gray
        # text_frame.alignment = PP_ALIGN.CENTER
        # text_frame.vertical_anchor = 4 
        
        try:
            slide.shapes.add_picture(visual_data['image'], left, top, width=width)
        except Exception as e:
            print("Error adding image:", e)
          
    def save(self, filename: str):
        self.prs.save(filename)